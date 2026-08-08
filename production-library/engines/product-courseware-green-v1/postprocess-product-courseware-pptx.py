#!/usr/bin/env python3
"""
商品培训课件 PPTX 后处理：
1) 注入企业统一中文字体（微软雅黑）到 theme + 全部 run
2) 将注意事项示例插图嵌入对应槽位

原因：@oai/artifact-tool 导出的 OOXML 只有字号/粗体，没有 typeface，
PowerPoint 会回退到主题 Calibri，中文全部变成同一默认字体。
"""
from __future__ import annotations

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 企业内训 PPT 统一字体（Windows 门店终端友好；macOS 有装亦可）
FONT_LATIN = "Microsoft YaHei"
FONT_EA = "微软雅黑"
FONT_CS = "Microsoft YaHei"

FONT_TAG = (
    f'<a:latin typeface="{FONT_LATIN}"/>'
    f'<a:ea typeface="{FONT_EA}"/>'
    f'<a:cs typeface="{FONT_CS}"/>'
)


def inject_fonts_into_rpr(xml: str) -> str:
    """在每个 a:rPr 内补 typeface；已有 typeface 的跳过。"""

    def repl(match: re.Match[str]) -> str:
        open_tag, body, close = match.group(1), match.group(2), match.group(3)
        if "typeface=" in body or "typeface=" in open_tag:
            return match.group(0)
        # self-closing <a:rPr .../>
        if close.startswith("/>") or open_tag.endswith("/>"):
            # convert self-closing to open+fonts+close
            open_clean = open_tag.rstrip()
            if open_clean.endswith("/>"):
                open_clean = open_clean[:-2] + ">"
            return f"{open_clean}{FONT_TAG}</a:rPr>"
        return f"{open_tag}{FONT_TAG}{body}{close}"

    # paired rPr
    xml = re.sub(
        r"(<a:rPr\b[^>]*>)(.*?)(</a:rPr>)",
        repl,
        xml,
        flags=re.DOTALL,
    )
    # remaining self-closing
    xml = re.sub(
        r"<a:rPr\b([^>]*?)/>",
        lambda m: (
            m.group(0)
            if "typeface=" in m.group(0)
            else f'<a:rPr{m.group(1)}>{FONT_TAG}</a:rPr>'
        ),
        xml,
    )
    return xml


def patch_theme_fonts(xml: str) -> str:
    """主题 major/minor 字体改为微软雅黑，避免默认 Calibri。"""

    def scheme_repl(match: re.Match[str]) -> str:
        kind = match.group(1)  # majorFont | minorFont
        return (
            f"<a:{kind}>"
            f'<a:latin typeface="{FONT_LATIN}"/>'
            f'<a:ea typeface="{FONT_EA}"/>'
            f'<a:cs typeface="{FONT_CS}"/>'
            f"</a:{kind}>"
        )

    xml = re.sub(
        r"<a:(majorFont|minorFont)>.*?</a:\1>",
        scheme_repl,
        xml,
        flags=re.DOTALL,
    )
    return xml


def fix_fonts_in_pptx(src: Path, dst: Path) -> dict:
    stats = {"slides": 0, "theme": 0, "other": 0}
    tmp = dst.with_suffix(".fontfix.zip")
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(
        tmp, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            name = info.filename
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                text = data.decode("utf-8")
                new = inject_fonts_into_rpr(text)
                data = new.encode("utf-8")
                stats["slides"] += 1
            elif name.startswith("ppt/slideMasters/") and name.endswith(".xml"):
                text = data.decode("utf-8")
                data = inject_fonts_into_rpr(text).encode("utf-8")
                stats["other"] += 1
            elif name.startswith("ppt/slideLayouts/") and name.endswith(".xml"):
                text = data.decode("utf-8")
                data = inject_fonts_into_rpr(text).encode("utf-8")
                stats["other"] += 1
            elif name == "ppt/theme/theme1.xml":
                text = data.decode("utf-8")
                data = patch_theme_fonts(text).encode("utf-8")
                stats["theme"] += 1
            zout.writestr(info, data)
    shutil.move(tmp, dst)
    return stats


def embed_precaution_images(pptx_path: Path, image_map: dict[str, Path]) -> int:
    """按 shape name 把注意事项槽位换成示例图。"""
    prs = Presentation(str(pptx_path))
    # 注意事项是最后一页
    slide = prs.slides[-1]
    placed = 0

    # 先收集目标框位置（surface 形状）
    targets: list[tuple[str, int, int, int, int]] = []
    for shape in slide.shapes:
        name = shape.name or ""
        m = re.match(r"precaution-asset-(\d+)-surface", name)
        if not m:
            continue
        idx = m.group(1)
        key = f"precaution-asset-{idx}"
        if key not in image_map:
            continue
        targets.append((key, shape.left, shape.top, shape.width, shape.height))

    for key, left, top, width, height in targets:
        img = image_map[key]
        if not img.exists():
            continue
        # 叠在原槽位上；保留卡片标题文字
        slide.shapes.add_picture(str(img), left, top, width=width, height=height)
        placed += 1

    prs.save(str(pptx_path))
    return placed


def main() -> int:
    if len(sys.argv) < 2:
        print(
            "usage: postprocess-product-courseware-pptx.py <pptx> [--assets-dir DIR]",
            file=sys.stderr,
        )
        return 2

    pptx = Path(sys.argv[1]).resolve()
    assets_dir = None
    if "--assets-dir" in sys.argv:
        assets_dir = Path(sys.argv[sys.argv.index("--assets-dir") + 1]).resolve()
    else:
        assets_dir = pptx.parent / "assets" / "precautions"

    work = pptx.with_suffix(".font-patched.pptx")
    stats = fix_fonts_in_pptx(pptx, work)
    shutil.move(work, pptx)

    image_map = {
        "precaution-asset-0": assets_dir / "01-not-replace-drug.png",
        "precaution-asset-1": assets_dir / "02-special-groups.png",
        "precaution-asset-2": assets_dir / "03-with-meal.png",
        "precaution-asset-3": assets_dir / "04-see-doctor.png",
    }
    placed = 0
    if assets_dir.exists():
        placed = embed_precaution_images(pptx, image_map)

    # 再跑一遍字体注入（python-pptx 另存可能不保留部分 run 字体）
    work2 = pptx.with_suffix(".font-patched2.pptx")
    stats2 = fix_fonts_in_pptx(pptx, work2)
    shutil.move(work2, pptx)

    print(
        {
            "pptx": str(pptx),
            "font": f"{FONT_EA} / {FONT_LATIN}",
            "font_patch": stats,
            "font_repatch": stats2,
            "images_placed": placed,
            "assets_dir": str(assets_dir),
        }
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
